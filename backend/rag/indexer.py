from __future__ import annotations

import logging
import re

from .embedder import OllamaEmbeddingFunction
from .store import delete_collection, get_or_create_collection

logger = logging.getLogger("simplechat.rag.indexer")

# matches markdown headings, numbered sections (4.2 Title), and ALL-CAPS lines
_HEADING_RE = re.compile(
    r'^(?:#{1,6}\s+|\d+(?:\.\d+)*[\s.]+[A-Z])',
)

def _is_heading(line: str) -> bool:
    s = line.strip()
    if not s or len(s) > 120:
        return False
    if _HEADING_RE.match(s):
        return True
    # ALL CAPS with at least two words (avoids "NOTE:", "TODO" etc.)
    return s.isupper() and len(s) >= 8 and " " in s


def _tail_sentences(text: str, n: int = 2) -> str:
    """Return the last n sentences of text for context-preserving overlap."""
    parts = [p for p in re.split(r'(?<=[.!?])\s+', text.strip()) if p]
    return " ".join(parts[-n:]) if parts else ""

_SUPPORTED_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/json",
    "application/pdf",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

DATASET_ALLOWED_MIME_TYPES = _SUPPORTED_MIME_TYPES


def extract_text(content: bytes, mime_type: str) -> str:
    """Convert file bytes to plain text for indexing."""
    if mime_type in ("text/plain", "text/markdown", "text/csv", "application/json"):
        return content.decode("utf-8", errors="replace")

    if mime_type == "application/pdf":
        import io

        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import io

        import docx as _docx
        doc = _docx.Document(io.BytesIO(content))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts)

    if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        import io

        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.values:
                parts.append("\t".join("" if v is None else str(v) for v in row))
        wb.close()
        return "\n".join(parts)

    if mime_type == "application/vnd.ms-excel":
        import xlrd
        wb = xlrd.open_workbook(file_contents=content)
        parts = []
        for name in wb.sheet_names():
            sheet = wb.sheet_by_name(name)
            for r in range(sheet.nrows):
                parts.append("\t".join(str(sheet.cell_value(r, c)) for c in range(sheet.ncols)))
        return "\n".join(parts)

    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        import io

        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.extend(
                        para.text.strip()
                        for para in shape.text_frame.paragraphs
                        if para.text.strip()
                    )
        return "\n".join(parts)

    return ""


def chunk_text(text: str, chunk_size: int = 512, overlap: int = 64) -> list[str]:
    """Split text into overlapping chunks with heading context and sentence-aware overlap.

    Each chunk is prefixed with the current section heading so retrieved chunks
    carry their document location. Heading is stored separately from content so
    every hard-split piece is also prefixed, not just the first.
    """
    if not text.strip():
        return []

    # parse into (is_heading, text) blocks
    blocks: list[tuple[bool, str]] = []
    current_lines: list[str] = []

    for line in text.splitlines():
        stripped = line.strip()
        if _is_heading(stripped):
            if current_lines:
                blocks.append((False, "\n".join(current_lines)))
                current_lines = []
            blocks.append((True, stripped))
        elif stripped:
            current_lines.append(stripped)
        else:
            if current_lines:
                blocks.append((False, "\n".join(current_lines)))
                current_lines = []
    if current_lines:
        blocks.append((False, "\n".join(current_lines)))

    # assemble (heading, body) pairs — heading stored separately for re-use during hard-split
    pending: list[tuple[str, str]] = []
    buf_body = ""
    buf_heading = ""

    for is_h, block_text in blocks:
        if is_h:
            if buf_body:
                pending.append((buf_heading, buf_body))
                buf_body = ""
            buf_heading = block_text
            continue

        prefix_len = len(f"[{buf_heading}]\n") if buf_heading else 0
        candidate = buf_body + ("\n\n" + block_text if buf_body else block_text)

        if prefix_len + len(candidate) <= chunk_size:
            buf_body = candidate
        else:
            if buf_body:
                pending.append((buf_heading, buf_body))
                tail = _tail_sentences(buf_body)
                buf_body = (tail + "\n\n" + block_text if tail else block_text).strip()
            else:
                buf_body = block_text  # single oversized block — hard-split below

    if buf_body:
        pending.append((buf_heading, buf_body))

    # emit chunks, re-applying heading prefix to every hard-split piece
    result: list[str] = []
    for h, body in pending:
        prefix = f"[{h}]\n" if h else ""
        full = (prefix + body).strip()
        if len(full) <= chunk_size:
            result.append(full)
        else:
            avail = max(overlap + 1, chunk_size - len(prefix))
            start = 0
            while start < len(body):
                result.append((prefix + body[start : start + avail]).strip())
                start += avail - overlap

    return [c for c in result if c.strip()]


def index_file(
    dataset_id: int,
    file_id: int,
    filename: str,
    content: bytes,
    mime_type: str,
    base_url: str,
    model: str,
) -> int:
    """Extract, chunk, and add a file's content to the dataset collection. Returns chunk count."""
    text = extract_text(content, mime_type)
    if not text.strip():
        logger.warning("index_file: no text extracted from %s (mime=%s)", filename, mime_type)
        return 0
    chunks = chunk_text(text)
    if not chunks:
        logger.warning("index_file: chunking produced 0 chunks for %s (%d chars)", filename, len(text))
        return 0
    logger.info("index_file: %s → %d chars → %d chunks (dataset=%d model=%s)", filename, len(text), len(chunks), dataset_id, model)
    embed_fn = OllamaEmbeddingFunction(base_url, model)
    col = get_or_create_collection(dataset_id, embed_fn)
    ids = [f"f{file_id}_c{i}" for i in range(len(chunks))]
    metadatas = [{"file_id": file_id, "filename": filename, "chunk": i} for i in range(len(chunks))]
    col.add(documents=chunks, ids=ids, metadatas=metadatas)
    logger.info("index_file: indexed %d chunks for file %d in dataset %d", len(chunks), file_id, dataset_id)
    return len(chunks)


def reindex_dataset(dataset_id: int, files: list, base_url: str, model: str) -> int:
    """Rebuild the collection from scratch using the supplied file list."""
    delete_collection(dataset_id)
    total = 0
    for f in files:
        total += index_file(dataset_id, f.id, f.filename, f.content, f.mime_type, base_url, model)
    return total
