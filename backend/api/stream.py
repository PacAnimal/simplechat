import asyncio
import base64
import json
import logging

import aiofiles
from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import StreamingResponse
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from .. import model_registry, sse_events
from ..auth import get_current_profile
from ..config import settings
from ..database import get_db
from ..event_logging import audit_message, log_event
from ..models import Attachment, Chat, Dataset, GeneratedImage, Message, Profile, utcnow
from ..providers import AnthropicProvider, OllamaProvider, OpenAIProvider
from ..providers.stub_provider import StubProvider
from ..schemas import SendMessageRequest
from .deps import get_owned_chat

logger = logging.getLogger(__name__)

_TEXT_MIME_TYPES = {"text/plain", "text/markdown", "text/csv", "application/json"}
_IMAGE_MIME_TYPES = {"image/png", "image/jpeg", "image/gif", "image/webp"}
_DOC_MIME_TYPES = {
    "application/pdf",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# per-chat lock: prevent double-sending to the same chat
_chat_locks: dict[int, asyncio.Lock] = {}


def _get_chat_lock(chat_id: int) -> asyncio.Lock:
    if chat_id not in _chat_locks:
        _chat_locks[chat_id] = asyncio.Lock()
    return _chat_locks[chat_id]


async def _attachment_text(att: Attachment) -> str:
    if att.mime_type in _DOC_MIME_TYPES:
        try:
            body = await asyncio.to_thread(_doc_to_text, att.path, att.mime_type)
            return f"\n\n[Attached document: {att.filename}]\n```\n{body}\n```"
        except Exception:
            logger.warning(
                "Failed to extract document %s (%s)", att.filename, att.path, exc_info=True
            )
            return ""
    if att.mime_type not in _TEXT_MIME_TYPES:
        return ""
    try:
        async with aiofiles.open(att.path, encoding="utf-8", errors="replace") as f:
            body = await f.read()
        return f"\n\n[Attached file: {att.filename}]\n```\n{body}\n```"
    except Exception:
        logger.warning(
            "Failed to read attachment %s (%s)", att.filename, att.path, exc_info=True
        )
        return ""


def _doc_to_text(path: str, mime: str) -> str:
    """Convert an office/PDF document to plain text for embedding in a message."""
    if mime == "application/pdf":
        from pypdf import PdfReader
        reader = PdfReader(path)
        parts = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"[Page {i}]")
                parts.append(text)
        return "\n".join(parts)

    if mime in ("application/vnd.ms-excel",):
        import xlrd
        wb = xlrd.open_workbook(path)
        parts = []
        for name in wb.sheet_names():
            sheet = wb.sheet_by_name(name)
            if sheet.nrows == 0:
                continue
            parts.append(f"[Sheet: {name}]")
            for r in range(sheet.nrows):
                parts.append("\t".join(str(sheet.cell_value(r, c)) for c in range(sheet.ncols)))
        return "\n".join(parts)

    if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = list(ws.values)
            if not rows:
                continue
            parts.append(f"[Sheet: {name}]")
            for row in rows:
                parts.append("\t".join("" if v is None else str(v) for v in row))
        wb.close()
        return "\n".join(parts)

    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx as _docx
        doc = _docx.Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts)

    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)

    return ""


async def _attachment_image_block(att: Attachment) -> dict | None:
    if att.mime_type not in _IMAGE_MIME_TYPES:
        return None
    try:
        async with aiofiles.open(att.path, "rb") as f:
            data = await f.read()
        return {"type": "image", "media_type": att.mime_type, "data": base64.b64encode(data).decode()}
    except Exception:
        logger.warning(
            "Failed to read image attachment %s (%s)", att.filename, att.path, exc_info=True
        )
        return None


router = APIRouter(prefix="/chats", tags=["stream"])


_PROVIDER_KEYS = {
    "openai": "openai_api_key",
    "anthropic": "anthropic_api_key",
    "ollama": "ollama_api_url",
}
_PROVIDER_LABELS = {"openai": "OpenAI", "anthropic": "Anthropic", "ollama": "Ollama"}


def _check_provider_configured(provider: str) -> None:
    """Raise 503 before the stream starts so the status code is still settable."""
    key_attr = _PROVIDER_KEYS.get(provider)
    if key_attr and not getattr(settings, key_attr, None):
        label = _PROVIDER_LABELS.get(provider, provider)
        raise HTTPException(
            status_code=503, detail=f"{label} is not configured on this server"
        )


def _get_provider(chat: Chat):
    if settings.stub_providers:
        return StubProvider(chat.provider)
    if chat.provider == "openai":
        return OpenAIProvider()
    if chat.provider == "anthropic":
        return AnthropicProvider()
    if chat.provider == "ollama":
        return OllamaProvider()
    raise HTTPException(status_code=400, detail=f"Unknown provider: {chat.provider}")


async def _build_messages(chat_id: int, db: AsyncSession) -> list[dict]:
    result = await db.execute(
        select(Message)
        .where(Message.chat_id == chat_id)
        .options(
            selectinload(Message.attachments),
            selectinload(Message.generated_images),
        )
        .order_by(Message.created_at)
    )
    rows = result.scalars().all()
    out = []
    for m in rows:
        if m.role == "user":
            text = m.content
            for att in m.attachments:
                text += await _attachment_text(att)
            image_blocks = [b for att in m.attachments if (b := await _attachment_image_block(att)) is not None]
            if image_blocks:
                content: str | list = [{"type": "text", "text": text}, *image_blocks]
            else:
                content = text
            out.append({"role": "user", "content": content})
        elif m.role == "assistant":
            content = m.content
            for img in m.generated_images:
                content += (
                    f"\n\n[Generated image — path: {img.path} | prompt: {img.prompt}]"
                )
            out.append({"role": "assistant", "content": content})
    return out


async def _save_user_message(
    chat_id: int, user_content: str, attachment_ids: list[int], db: AsyncSession
) -> Message:
    user_msg = Message(chat_id=chat_id, role="user", content=user_content)
    db.add(user_msg)
    await db.commit()
    await db.refresh(user_msg)
    if attachment_ids:
        result = await db.execute(
            select(Attachment).where(
                Attachment.id.in_(attachment_ids),
                Attachment.chat_id == chat_id,
                Attachment.message_id.is_(None),
            )
        )
        for att in result.scalars().all():
            att.message_id = user_msg.id
        await db.commit()
    return user_msg


async def _save_assistant_message(
    chat_id: int,
    content: str,
    thinking: str | None,
    generated_images: list[dict],
    chat: Chat,
    db: AsyncSession,
) -> Message:
    assistant_msg = Message(
        chat_id=chat_id, role="assistant", content=content, thinking=thinking
    )
    db.add(assistant_msg)
    chat.updated_at = utcnow()
    await db.commit()
    await db.refresh(assistant_msg)
    for img in generated_images:
        gi = GeneratedImage(
            chat_id=chat_id,
            message_id=assistant_msg.id,
            prompt=img.get("prompt", ""),
            path=img.get("path", ""),
        )
        db.add(gi)
    if generated_images:
        await db.commit()
    return assistant_msg


async def _generate_chat_title(
    user_content: str, assistant_content: str, provider: str, model: str
) -> str:
    """Use a lightweight LLM call to produce a short chat title (5 words max)."""
    if settings.stub_providers:
        return user_content[:60].strip() or "New Chat"
    system = "Generate a short title for the conversation shown. Reply with only the title — 5 words max, no quotes, no punctuation."
    user_msg = f"User: {user_content[:400]}\nAssistant: {assistant_content[:400]}"
    try:
        if provider == "anthropic" and settings.anthropic_api_key:
            import anthropic

            client = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)
            resp = await client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=20,
                system=system,
                messages=[{"role": "user", "content": user_msg}],
            )
            return resp.content[0].text.strip()[:100]
        if provider == "openai" and settings.openai_api_key:
            from openai import AsyncOpenAI

            client = AsyncOpenAI(api_key=settings.openai_api_key)
            resp = await client.chat.completions.create(
                model="gpt-4o-mini",
                max_tokens=20,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": user_msg},
                ],
            )
            return (resp.choices[0].message.content or "").strip()[:100]
        if provider == "ollama" and settings.ollama_api_url:
            from openai import AsyncOpenAI

            base_url = settings.ollama_api_url.rstrip("/") + "/v1"
            client = AsyncOpenAI(api_key="ollama", base_url=base_url)
            resp = await client.chat.completions.create(
                model=model,
                max_tokens=20,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": user_msg},
                ],
            )
            return (resp.choices[0].message.content or "").strip()[:100]
    except Exception:
        logger.exception("Title generation failed")
    return user_content[:60].strip() or "New Chat"


_SUMMARY_CHAR_THRESHOLD = 6_000  # ~1500 tokens; trigger point for small-model context limits
_SUMMARY_KEEP_RECENT = 8         # messages kept verbatim (4 user+assistant turns)


async def _summarize_history(
    messages: list[dict], model: str, ollama_url: str
) -> list[dict]:
    """Compress old messages into a summary when the conversation grows too long.

    Keeps the most recent _SUMMARY_KEEP_RECENT messages verbatim so the model
    retains immediate context; older turns are distilled into a single summary
    exchange at the front of the list.
    """
    total = sum(len(str(m.get("content", ""))) for m in messages)
    if total <= _SUMMARY_CHAR_THRESHOLD or len(messages) <= _SUMMARY_KEEP_RECENT:
        return messages

    old = messages[:-_SUMMARY_KEEP_RECENT]
    recent = messages[-_SUMMARY_KEEP_RECENT:]

    if len(old) < 3:
        return messages

    history_text = "\n".join(
        f"{m['role'].capitalize()}: {str(m.get('content', ''))[:400]}"
        for m in old
    )
    try:
        from openai import AsyncOpenAI

        client = AsyncOpenAI(api_key="ollama", base_url=ollama_url.rstrip("/") + "/v1")
        resp = await client.chat.completions.create(
            model=model,
            max_tokens=200,
            messages=[
                {
                    "role": "system",
                    "content": "Summarize the following conversation in 2-3 sentences, preserving key facts and decisions.",
                },
                {"role": "user", "content": history_text},
            ],
        )
        summary = (resp.choices[0].message.content or "").strip()
        if summary:
            return [
                {"role": "user", "content": f"[Summary of earlier conversation: {summary}]"},
                {"role": "assistant", "content": "Understood."},
                *recent,
            ]
    except Exception:
        logger.warning("Conversation summarization failed", exc_info=True)

    # fallback: just drop old messages rather than blowing the context window
    return recent


async def _reformulate_rag_query(
    history: list[dict], user_query: str, ollama_url: str, model: str
) -> str:
    """Rewrite user query as a standalone retrieval-optimised search string.

    Resolves pronouns and follow-up references using the conversation history so
    queries like "expand on the second point" actually retrieve something useful.
    Falls back to the original query on any error.
    """
    if not history:
        return user_query
    history_text = "\n".join(
        f"{m['role'].capitalize()}: {str(m.get('content', ''))[:300]}"
        for m in history[-4:]
        if isinstance(m.get("content"), str)
    )
    prompt = (
        f"Conversation so far:\n{history_text}\n\n"
        f"Latest user message: {user_query}\n\n"
        "Rewrite the latest user message as a concise, standalone search query for a "
        "document retrieval system. Resolve any pronouns or references using the "
        "conversation context. Return only the rewritten query, nothing else."
    )
    try:
        from openai import AsyncOpenAI

        client = AsyncOpenAI(api_key="ollama", base_url=ollama_url.rstrip("/") + "/v1")
        resp = await client.chat.completions.create(
            model=model,
            max_tokens=80,
            temperature=0,
            messages=[
                {"role": "system", "content": "You are a search query optimizer."},
                {"role": "user", "content": prompt},
            ],
        )
        reformulated = (resp.choices[0].message.content or "").strip()
        if reformulated:
            logger.info("RAG: reformulated query %r → %r", user_query[:80], reformulated[:80])
            return reformulated
    except Exception:
        logger.debug("RAG: query reformulation failed, using original", exc_info=True)
    return user_query


async def _event_stream(
    chat: Chat,
    user_content: str,
    attachment_ids: list[int],
    db: AsyncSession,
    request: Request,
    lock: asyncio.Lock,
    chat_id: int,
):
    lock_released = False

    def _release_lock():
        nonlocal lock_released
        if not lock_released:
            lock_released = True
            lock.release()

    try:
        user_msg = await _save_user_message(chat.id, user_content, attachment_ids, db)
        yield f"data: {json.dumps({'type': sse_events.USER_MESSAGE_SAVED, 'message_id': user_msg.id})}\n\n"
        messages = await _build_messages(chat.id, db)
        real_model = model_registry.resolve_model_id(chat.provider, chat.model)

        # summarize old turns before injecting RAG so context window stays manageable
        if chat.provider == "ollama" and settings.ollama_api_url:
            messages = await _summarize_history(messages, real_model, settings.ollama_api_url)

        # inject RAG context for ollama chats with an active dataset
        if chat.provider == "ollama" and chat.dataset_id and settings.ollama_api_url:
            dataset = await db.get(Dataset, chat.dataset_id)
            if dataset and messages:
                try:
                    from ..rag.embedder import EMBED_MODEL, OllamaEmbeddingFunction
                    from ..rag.store import query_collection
                    embed_fn = OllamaEmbeddingFunction(settings.ollama_api_url, EMBED_MODEL)
                    # reformulate query using conversation history to resolve references and pronouns
                    history_query = await _reformulate_rag_query(
                        messages[:-1], user_content, settings.ollama_api_url, real_model
                    )
                    logger.info(
                        "RAG: chat=%d dataset=%d(%r) query=%r",
                        chat.id, chat.dataset_id, dataset.name, history_query[:120],
                    )
                    chunks = await asyncio.to_thread(
                        query_collection, chat.dataset_id, embed_fn, history_query
                    )
                    if chunks:
                        rag_text = "\n\n---\n\n".join(chunks)
                        prefix = f"[Relevant context from dataset \"{dataset.name}\"]\n{rag_text}\n\n[User message]\n"
                        last = messages[-1]
                        if isinstance(last["content"], str):
                            last["content"] = prefix + last["content"]
                        elif isinstance(last["content"], list):
                            for block in last["content"]:
                                if block.get("type") == "text":
                                    block["text"] = prefix + block["text"]
                                    break
                        logger.info("RAG: injected %d chunk(s) into chat %d message", len(chunks), chat.id)
                    else:
                        logger.warning("RAG: 0 chunks returned for chat %d — model will have no context", chat.id)
                except Exception:
                    logger.warning("RAG context injection failed for chat %d", chat.id, exc_info=True)
        elif chat.provider == "ollama" and not chat.dataset_id:
            logger.debug("RAG: skipped for chat %d — no dataset selected", chat.id)

        provider = _get_provider(chat)

        full_content = ""
        thinking_content = ""
        generated_images: list[dict] = []

        async def _watch_disconnect():
            while not await request.is_disconnected():
                await asyncio.sleep(0.3)

        disconnect_task = asyncio.create_task(_watch_disconnect())
        try:
            stream = provider.stream_chat(messages, real_model, chat.web_search_enabled)
            async for event in stream:
                if disconnect_task.done():
                    logger.info("Client disconnected mid-stream for chat %d", chat.id)
                    return
                etype = event.get("type")
                if etype == sse_events.TEXT_DELTA:
                    full_content += event.get("content", "")
                elif etype == sse_events.THINKING_DELTA:
                    thinking_content += event.get("content", "")
                elif etype == sse_events.IMAGE_GENERATED:
                    generated_images.append(event)
                yield f"data: {json.dumps(event)}\n\n"
        finally:
            disconnect_task.cancel()

        assistant_msg = await _save_assistant_message(
            chat.id, full_content, thinking_content or None, generated_images, chat, db
        )

        # release lock here so the user can send another message while title generates
        yield f"data: {json.dumps({'type': sse_events.DONE, 'message_id': assistant_msg.id})}\n\n"
        _release_lock()

        if chat.title_is_default:
            title = await _generate_chat_title(
                user_content, full_content, chat.provider, real_model
            )
            chat.title = title
            chat.title_is_default = False
            await db.commit()
            yield f"data: {json.dumps({'type': sse_events.CHAT_TITLE, 'title': title})}\n\n"

    except Exception as e:
        logger.exception("Stream error for chat %d", chat.id)
        yield f"data: {json.dumps({'type': sse_events.ERROR, 'message': str(e)})}\n\n"
    finally:
        _release_lock()


@router.post("/{chat_id}/messages")
async def send_message(
    request: Request,
    chat_id: int,
    body: SendMessageRequest,
    profile: Profile = Depends(get_current_profile),
    db: AsyncSession = Depends(get_db),
):
    # verify ownership before acquiring lock — prevents lock leak on 404
    chat = await get_owned_chat(chat_id, profile.id, db)

    if not settings.stub_providers:
        _check_provider_configured(chat.provider)

    lock = _get_chat_lock(chat_id)
    if lock.locked():
        raise HTTPException(429, "A message is already being processed for this chat")
    await lock.acquire()

    log_event(profile.name, "send_message", chat_id=chat_id)
    if settings.audit_log:
        audit_message(profile.name, chat_id, body.content, provider=chat.provider, model=chat.model)

    return StreamingResponse(
        _event_stream(
            chat, body.content, body.attachment_ids, db, request, lock, chat_id
        ),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )
